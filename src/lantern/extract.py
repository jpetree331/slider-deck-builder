"""Attachment extraction — PDF / DOCX / PPTX in; plain text + key embedded
images out.

Framework-free on purpose (no FastAPI imports) — exercised headless by
scripts/verify_extract.py.

Design stance: extract-and-discard. Text lands in the source-notes textarea,
visible and editable (Sacred Invariant 6); images are shown as chips and ride
along to the ONE Haiku outline call so it can see the source material's
visual character — then both are dropped. Deck folders stay pictures + one
JSON; uploaded files are never written to disk.

Honest limitation (say it in the UI too): we pull the images embedded in the
document, not a rendering of its pages — rasterizing a PPTX layout needs
PowerPoint itself. Text + embedded imagery is what "revamp this deck" gets.
"""
import base64
import hashlib
import io
import logging

from PIL import Image

logger = logging.getLogger("lantern.extract")

MAX_FILE_BYTES = 20 * 1024 * 1024  # refuse uploads over 20 MB
MAX_TEXT_CHARS = 80_000  # keep the outline prompt (and its cost) sane
TRUNCATION_NOTE = ("\n\n[attachment truncated here — it was longer than "
                   "Lantern sends to the outline model]")
MAX_IMAGES = 8          # plenty for style; pennies of Haiku vision
MIN_IMAGE_BYTES = 3072  # skip bullet dots and icons
MIN_IMAGE_SIDE = 64
MAX_IMAGE_SIDE = 1024   # downscale before it ever leaves the server

KINDS = ("pdf", "docx", "pptx")


class ExtractError(Exception):
    """User-facing extraction failure — maps to HTTP 422."""


def extract_content(filename: str, data: bytes) -> dict:
    ext = (filename or "").rpartition(".")[2].lower()
    if ext not in KINDS:
        raise ExtractError(f"unsupported file type {'.' + ext if ext else ''!r}"
                           " — attach a .pdf, .docx, or .pptx")
    text, raw_images = {"pdf": _pdf, "docx": _docx, "pptx": _pptx}[ext](data)
    text = text.strip()
    images = _prepare_images(raw_images)
    if not text and not images:
        raise ExtractError(f"no readable text or images found in {filename} — "
                           "if it's a scanned PDF, paste the text instead")
    truncated = len(text) > MAX_TEXT_CHARS
    if truncated:
        text = text[:MAX_TEXT_CHARS] + TRUNCATION_NOTE
    logger.info("extracted %d chars + %d image(s) from %s (%s%s)",
                len(text), len(images), filename, ext,
                ", truncated" if truncated else "")
    return {"filename": filename, "kind": ext, "text": text,
            "chars": len(text), "truncated": truncated, "images": images}


def _prepare_images(raw: list) -> list:
    """(bytes, note) pairs -> deduped, size-filtered, downscaled JPEG b64.
    Prefers the largest originals when over the cap."""
    seen, candidates = set(), []
    for blob, note in raw:
        if not blob or len(blob) < MIN_IMAGE_BYTES:
            continue
        digest = hashlib.sha1(blob).digest()
        if digest in seen:
            continue
        seen.add(digest)
        candidates.append((blob, note))
    candidates.sort(key=lambda pair: len(pair[0]), reverse=True)
    images = []
    for blob, note in candidates:
        if len(images) >= MAX_IMAGES:
            break
        try:
            img = Image.open(io.BytesIO(blob))
            img.load()
        except Exception:
            continue  # unsupported/broken media — skip quietly
        if img.width < MIN_IMAGE_SIDE or img.height < MIN_IMAGE_SIDE:
            continue
        img.thumbnail((MAX_IMAGE_SIDE, MAX_IMAGE_SIDE))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        out = io.BytesIO()
        img.save(out, "JPEG", quality=80)
        images.append({
            "media_type": "image/jpeg",
            "data": base64.b64encode(out.getvalue()).decode("ascii"),
            "note": note,
        })
    return images


# ── per-format readers: return (text, [(image_bytes, note), ...]) ───────────

def _pdf(data: bytes):
    from pypdf import PdfReader
    from pypdf.errors import PyPdfError
    try:
        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                raise ExtractError("this PDF is password-protected — remove "
                                   "the password or paste the text")
        pages, images = [], []
        for i, page in enumerate(reader.pages, 1):
            pages.append((page.extract_text() or "").strip())
            try:
                for image_file in page.images:
                    images.append((image_file.data, f"page {i}"))
            except Exception:
                pass  # malformed embedded object — text still counts
    except ExtractError:
        raise
    except PyPdfError as e:
        raise ExtractError(f"couldn't read this PDF ({e}) — it may be "
                           "damaged; paste the text instead")
    return "\n\n".join(t for t in pages if t), images


def _docx(data: bytes):
    import docx
    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as e:  # python-docx raises bare ValueError/KeyError etc.
        raise ExtractError(f"couldn't read this Word file ({e}) — is it a "
                           "real .docx?")
    parts = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    images = []
    for part in document.part.related_parts.values():
        if getattr(part, "content_type", "").startswith("image/"):
            images.append((part.blob, "embedded in document"))
    return "\n".join(parts), images


def _pptx(data: bytes):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ExtractError(f"couldn't read this PowerPoint file ({e}) — is it "
                           "a real .pptx?")
    blocks, images = [], []
    for i, slide in enumerate(presentation.slides, 1):
        lines = [shape.text_frame.text.strip() for shape in slide.shapes
                 if shape.has_text_frame and shape.text_frame.text.strip()]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append((shape.image.blob, f"slide {i}"))
                except Exception:
                    pass  # linked-not-embedded pictures have no blob
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None and notes.text.strip():
                lines.append(f"Speaker notes: {notes.text.strip()}")
        if lines:
            blocks.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(blocks), images
