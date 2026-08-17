"""Deck exports — PPTX / PDF / ZIP, rebuilt on EVERY call from the slide PNGs.

Exports are derived artifacts, never a second source of truth (Sacred
Invariant 7). All paths come from store helpers, never hand-built.
"""
import json
import logging
import os
import re
import zipfile
from pathlib import Path

import img2pdf
from pptx import Presentation
from pptx.util import Inches

from . import store

logger = logging.getLogger("lantern.export")

FORMATS = ("pptx", "pdf", "zip")

# full-bleed 16:9 (Locked Decision 9)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


class NotFullyRendered(Exception):
    """Deck has unrendered slides and allow_partial wasn't set — HTTP 409."""


class NothingToExport(Exception):
    """Zero rendered slides — HTTP 409."""


def _slug(title: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", title.lower()).strip("-")
    return slug or "deck"


def export_deck(deck_id: str, fmt: str, allow_partial: bool = False) -> Path:
    if fmt not in FORMATS:
        raise ValueError(f"unknown export format {fmt!r}")
    deck = store.load_deck(deck_id)
    rendered = [s for s in deck["slides"]
                if s["render"] and s["render"]["status"] == "done"
                and store.find_slide_image(deck_id, s["n"]) is not None]
    missing = len(deck["slides"]) - len(rendered)
    if missing and not allow_partial:
        raise NotFullyRendered(
            f"{missing} slide(s) not rendered yet — render the deck first, "
            "or export with allow_partial=true to skip them")
    if not rendered:
        raise NothingToExport("no rendered slides to export")

    out_dir = store.exports_dir(deck_id)
    out_dir.mkdir(parents=True, exist_ok=True)
    out = out_dir / f"lantern-{_slug(deck['title'])}.{fmt}"
    tmp = out.with_name(out.name + ".tmp")
    paths = [store.find_slide_image(deck_id, s["n"]) for s in rendered]  # slides[].n order

    if fmt == "pptx":
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        blank = prs.slide_layouts[6]
        for path in paths:
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(path), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)
        prs.core_properties.title = deck["title"]
        prs.save(str(tmp))
    elif fmt == "pdf":
        # one page per slide; page size matches the image aspect
        tmp.write_bytes(img2pdf.convert([str(p) for p in paths]))
    else:  # zip: the PNGs + deck.json
        with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zf:
            for path in paths:
                zf.write(path, path.name)
            zf.write(store.deck_dir(deck_id) / "deck.json", "deck.json")

    os.replace(tmp, out)
    logger.info("exported deck %s as %s (%d slide(s)) -> %s",
                deck_id, fmt, len(paths), out.name)
    return out
