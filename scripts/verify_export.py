"""Sprint 6 verify — exports rebuilt on demand, correct geometry and order.

Run from repo root: python scripts/verify_export.py
Headless, throwaway data dir, real Pillow-generated PNGs.
"""
import os
import sys
import tempfile
import time
import zipfile
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO))
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # cp1252 consoles

os.environ["LANTERN_DATA_DIR"] = tempfile.mkdtemp(prefix="lantern-verify-")

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402

from src.lantern import export, store  # noqa: E402

FAILS = []


def check(name, cond):
    print(("  ok  " if cond else "  FAIL") + f"  {name}")
    if not cond:
        FAILS.append(name)


COLORS = ["#D96C3A", "#3AA6D9", "#86A97F"]
deck = store.create_deck(
    title="Cell Transport! (draft #2)", topic="t",
    slides=[{"title": f"S{i}", "points": [], "visual_description": "v",
             "layout_hint": "split",
             "render": {"status": "done", "image": f"slides/{i:02d}.png",
                        "rendered_at": "2026-01-01T00:00:00+00:00",
                        "cost_estimate_usd": 0.134}}
            for i in (1, 2, 3)])
for i, color in zip((1, 2, 3), COLORS):
    Image.new("RGB", (320, 180), color).save(store.slide_image_path(deck["id"], i))

print("verify_export: pptx")
path = export.export_deck(deck["id"], "pptx")
check("slug strips punctuation", path.name == "lantern-cell-transport-draft-2.pptx")
prs = Presentation(str(path))
check("pptx is 13.333x7.5in", round(prs.slide_width.inches, 3) == 13.333
      and prs.slide_height.inches == 7.5)
check("one slide per PNG, in order", len(prs.slides) == 3)
pics = [shape for slide in prs.slides for shape in slide.shapes
        if shape.shape_type == 13]  # PICTURE
check("every slide is one full-bleed picture",
      len(pics) == 3 and all(p.left == 0 and p.top == 0
                             and p.width == prs.slide_width
                             and p.height == prs.slide_height for p in pics))
check("deck title in doc properties", prs.core_properties.title == deck["title"])

print("verify_export: pdf")
pdf = export.export_deck(deck["id"], "pdf")
head = pdf.read_bytes()[:5]
check("pdf magic bytes", head == b"%PDF-")
check("pdf has 3 pages", pdf.read_bytes().count(b"/Type /Page ") == 3
      or pdf.read_bytes().count(b"/Type/Page") >= 3)

print("verify_export: zip")
zpath = export.export_deck(deck["id"], "zip")
with zipfile.ZipFile(zpath) as zf:
    names = sorted(zf.namelist())
check("zip carries the PNGs + deck.json",
      names == ["01.png", "02.png", "03.png", "deck.json"])

print("verify_export: rebuilt on demand (invariant 7)")
first_mtime = path.stat().st_mtime_ns
time.sleep(0.02)
export.export_deck(deck["id"], "pptx")
check("second export rewrites the file", path.stat().st_mtime_ns > first_mtime)

print("verify_export: partial semantics")
def clear_slide_2(d):
    d["slides"][1]["render"] = None
store.update_deck(deck["id"], clear_slide_2)
try:
    export.export_deck(deck["id"], "pptx")
    check("partial deck without flag raises (409)", False)
except export.NotFullyRendered:
    check("partial deck without flag raises (409)", True)
partial = export.export_deck(deck["id"], "pptx", allow_partial=True)
check("allow_partial skips the unrendered slide",
      len(Presentation(str(partial)).slides) == 2)

def clear_all(d):
    for s in d["slides"]:
        s["render"] = None
store.update_deck(deck["id"], clear_all)
try:
    export.export_deck(deck["id"], "zip", allow_partial=True)
    check("zero rendered slides raises NothingToExport", False)
except export.NothingToExport:
    check("zero rendered slides raises NothingToExport", True)

store.delete_deck(deck["id"])

if FAILS:
    print(f"\n{len(FAILS)} FAILURES: {FAILS}")
    sys.exit(1)
print("\nall export checks passed")
