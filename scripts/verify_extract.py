"""Attachment-extraction verify — real PPTX/DOCX/PDF fixtures built in-memory,
image plumbing to the outline call, HTTP endpoint contract. Fully offline.

Run from repo root: python scripts/verify_extract.py
"""
import base64
import io
import json
import os
import sys
import tempfile
from pathlib import Path
from types import SimpleNamespace

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO))
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")  # cp1252 consoles

os.environ["LANTERN_DATA_DIR"] = tempfile.mkdtemp(prefix="lantern-verify-")
os.environ["CHALK_DB_PATH"] = os.path.join(
    os.environ["LANTERN_DATA_DIR"], "chalk.db")

from PIL import Image  # noqa: E402

from src.lantern import extract  # noqa: E402
from src.lantern.outline import generate_outline  # noqa: E402

FAILS = []


def check(name, cond):
    print(("  ok  " if cond else "  FAIL") + f"  {name}")
    if not cond:
        FAILS.append(name)


import hashlib  # noqa: E402


def make_png(width=320, height=200, seed=1) -> bytes:
    """Truly incompressible PNG (hash-chain noise) so fixtures clear the
    icon-size filter like real photos do. Deterministic per seed."""
    need = width * height * 3
    out, h = bytearray(), bytes([seed])
    while len(out) < need:
        h = hashlib.sha256(h).digest()
        out += h
    buf = io.BytesIO()
    Image.frombytes("RGB", (width, height), bytes(out[:need])).save(buf, "PNG")
    return buf.getvalue()


def make_pdf_with_text(text: str) -> bytes:
    """Minimal valid one-page PDF with a real text object."""
    content = f"BT /F1 24 Tf 72 700 Td ({text}) Tj ET".encode()
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objs, 1):
        offsets.append(out.tell())
        out.write(b"%d 0 obj\n" % i + obj + b"\nendobj\n")
    xref = out.tell()
    out.write(b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1))
    for off in offsets:
        out.write(b"%010d 00000 n \n" % off)
    out.write(b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF"
              % (len(objs) + 1, xref))
    return out.getvalue()


print("verify_extract: pptx (text + notes + embedded picture)")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Cell Transport Basics"
slide.notes_slide.notes_text_frame.text = "Remind them about the lab Friday"
slide.shapes.add_picture(io.BytesIO(make_png()), Inches(1), Inches(2),
                         width=Inches(4))
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
slide2.shapes.title.text = "Osmosis in One Picture"
buf = io.BytesIO()
prs.save(buf)
result = extract.extract_content("old-deck.pptx", buf.getvalue())
check("slide text extracted with slide markers",
      "[Slide 1]" in result["text"] and "Cell Transport Basics" in result["text"]
      and "[Slide 2]" in result["text"])
check("speaker notes captured", "lab Friday" in result["text"])
check("embedded picture extracted", len(result["images"]) == 1
      and result["images"][0]["note"] == "slide 1")
check("image re-encoded as jpeg b64", result["images"][0]["media_type"] == "image/jpeg"
      and len(base64.b64decode(result["images"][0]["data"])) > 500)

print("verify_extract: docx (paragraphs + table + picture)")
import docx  # noqa: E402

doc = docx.Document()
doc.add_paragraph("Unit goals for the semester")
table = doc.add_table(rows=1, cols=2)
table.rows[0].cells[0].text = "Week 1"
table.rows[0].cells[1].text = "Membranes"
doc.add_picture(io.BytesIO(make_png(400, 300, seed=2)), width=Inches(2))
buf = io.BytesIO()
doc.save(buf)
result = extract.extract_content("notes.docx", buf.getvalue())
check("paragraphs extracted", "Unit goals" in result["text"])
check("table rows flattened", "Week 1 | Membranes" in result["text"])
check("docx picture extracted", len(result["images"]) == 1)

print("verify_extract: pdf")
result = extract.extract_content("handout.pdf",
                                 make_pdf_with_text("Membrane transport basics"))
check("pdf text extracted", "Membrane transport basics" in result["text"])

import img2pdf  # noqa: E402
image_only = img2pdf.convert([make_png(800, 450)])
result = extract.extract_content("scan.pdf", image_only)
check("image-only pdf yields its images instead of erroring",
      result["text"] == "" and len(result["images"]) == 1)

print("verify_extract: guardrails")
try:
    extract.extract_content("virus.exe", b"MZ")
    check("unknown extension raises", False)
except extract.ExtractError:
    check("unknown extension raises", True)
try:
    extract.extract_content("broken.pdf", b"not a pdf at all")
    check("damaged pdf raises cleanly", False)
except extract.ExtractError:
    check("damaged pdf raises cleanly", True)

real_cap = extract.MAX_TEXT_CHARS
extract.MAX_TEXT_CHARS = 10
result = extract.extract_content("handout.pdf",
                                 make_pdf_with_text("Membrane transport basics and much more"))
check("over-budget text truncated with note",
      result["truncated"] and "truncated" in result["text"])
extract.MAX_TEXT_CHARS = real_cap

tiny = make_png(16, 16)  # icon-sized — must be filtered out
prs2 = Presentation()
s = prs2.slides.add_slide(prs2.slide_layouts[5])
s.shapes.title.text = "T"
s.shapes.add_picture(io.BytesIO(tiny), Inches(1), Inches(1))
big = make_png(2400, 1350)
s.shapes.add_picture(io.BytesIO(big), Inches(2), Inches(2), width=Inches(3))
s.shapes.add_picture(io.BytesIO(big), Inches(3), Inches(3), width=Inches(3))
buf = io.BytesIO()
prs2.save(buf)
result = extract.extract_content("mixed.pptx", buf.getvalue())
check("tiny icons filtered, duplicates deduped", len(result["images"]) == 1)
w, h = Image.open(io.BytesIO(base64.b64decode(result["images"][0]["data"]))).size
check("oversized image downscaled to <=1024", max(w, h) <= 1024)

print("verify_extract: outline vision plumbing (stub client)")
VALID = json.dumps({
    "title": "T",
    "style_guide": {"palette": ["#101418", "#F2E9DC", "#3AA6D9"],
                    "typography": "serif", "motif": "m",
                    "art_direction": "one cohesive paragraph", "tone": "warm"},
    "slides": [{"title": "S", "points": [], "layout_hint": "title card",
                "visual_description": "v"}],
})
captured = {}


class StubClient:
    def __init__(self):
        self.messages = SimpleNamespace(create=self._create)

    def _create(self, **kwargs):
        captured.update(kwargs)
        return SimpleNamespace(content=[SimpleNamespace(type="text", text=VALID)])


imgs = [{"media_type": "image/jpeg", "data": "aGVsbG8=", "note": "slide 1"}]
generate_outline("revamp this deck", client=StubClient(), images=imgs)
content = captured["messages"][0]["content"]
check("image block precedes the text block",
      isinstance(content, list) and content[0]["type"] == "image"
      and content[0]["source"]["media_type"] == "image/jpeg")
check("text block flags the attached visuals",
      content[-1]["type"] == "text" and "ATTACHED VISUALS" in content[-1]["text"])
generate_outline("plain", client=StubClient())
check("no images -> plain string content",
      isinstance(captured["messages"][0]["content"], str))

print("verify_extract: HTTP endpoint")
import asyncio  # noqa: E402

import httpx  # noqa: E402

from src.lantern.api import app  # noqa: E402


async def http_checks():
    transport = httpx.ASGITransport(app=app)
    async with httpx.AsyncClient(transport=transport,
                                 base_url="http://test", timeout=30) as client:
        prs3 = Presentation()
        s3 = prs3.slides.add_slide(prs3.slide_layouts[5])
        s3.shapes.title.text = "Wire Test"
        s3.shapes.add_picture(io.BytesIO(make_png()), Inches(1), Inches(1))
        out = io.BytesIO()
        prs3.save(out)
        r = await client.post("/api/extract", files={
            "file": ("wire.pptx", out.getvalue(),
                     "application/vnd.openxmlformats-officedocument"
                     ".presentationml.presentation")})
        body = r.json()
        check("endpoint extracts over the wire",
              r.status_code == 200 and "Wire Test" in body["text"]
              and len(body["images"]) == 1)
        r = await client.post("/api/extract",
                              files={"file": ("x.txt", b"hi", "text/plain")})
        check("unsupported type 422s with a friendly message",
              r.status_code == 422 and ".pdf" in r.json()["detail"])
        real_max = extract.MAX_FILE_BYTES
        extract.MAX_FILE_BYTES = 10
        r = await client.post("/api/extract",
                              files={"file": ("big.pdf", b"x" * 100,
                                              "application/pdf")})
        check("oversize upload 413s", r.status_code == 413)
        extract.MAX_FILE_BYTES = real_max


asyncio.run(http_checks())

print("verify_extract: live vision outline (real Haiku, ~1 cent)")
if os.environ.get("ANTHROPIC_API_KEY"):
    live_img = {"media_type": "image/jpeg", "note": "slide 1",
                "data": base64.b64encode(
                    (lambda b: (Image.open(io.BytesIO(make_png(320, 180)))
                                .convert("RGB").save(b, "JPEG"), b.getvalue())[1]
                     )(io.BytesIO())).decode("ascii")}
    live = generate_outline(
        "Two slides about cell membranes, matching the attached image's colors",
        slide_count_hint=2, images=[live_img])
    check("live vision outline validates", len(live.slides) == 2
          and live.style_guide.art_direction.strip() != "")
else:
    print("  SKIPPED — ANTHROPIC_API_KEY not set")

if FAILS:
    print(f"\n{len(FAILS)} FAILURES: {FAILS}")
    sys.exit(1)
print("\nall extract checks passed")
